"""Generate small, synthetic attachment fixtures for offline QA only."""

from __future__ import annotations

import argparse
import math
import struct
import wave
import zipfile
import zlib
from pathlib import Path


MARKERS = {
    "txt": "TXT_OFFLINE_CONTENT_5412",
    "py": "CODE_OFFLINE_CONTENT_8841",
    "pdf": "PDF_OFFLINE_CONTENT_1937",
    "docx": "DOCX_OFFLINE_CONTENT_2764",
    "xlsx": "XLSX_OFFLINE_CONTENT_6285",
    "pptx": "PPTX_OFFLINE_CONTENT_7359",
    "png": "PNG_OFFLINE_CONTENT_4028",
    "wav": "WAV_OFFLINE_CONTENT_9516",
    "mp4": "VIDEO_OFFLINE_CONTENT_3175",
    "zip": "ZIP_OFFLINE_CONTENT_8603",
}


def _write_pdf(path: Path, marker: str) -> None:
    stream = f"BT /F1 18 Tf 72 720 Td ({marker}) Tj ET".encode("ascii")
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Resources << /Font << /F1 5 0 R >> >> /Contents 4 0 R >>"
        ),
        b"<< /Length " + str(len(stream)).encode("ascii") + b" >>\nstream\n"
        + stream
        + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    payload = bytearray(b"%PDF-1.4\n%synthetic\n")
    offsets = [0]
    for index, body in enumerate(objects, start=1):
        offsets.append(len(payload))
        payload.extend(f"{index} 0 obj\n".encode("ascii"))
        payload.extend(body)
        payload.extend(b"\nendobj\n")
    xref = len(payload)
    payload.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    payload.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        payload.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    payload.extend(
        (
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref}\n%%EOF\n"
        ).encode("ascii")
    )
    path.write_bytes(payload)


def _write_docx(path: Path, marker: str) -> None:
    content_types = """<?xml version="1.0" encoding="UTF-8"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>"""
    relationships = """<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>"""
    document = f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:body><w:p><w:r><w:t>{marker}</w:t></w:r></w:p><w:sectPr/></w:body>
</w:document>"""
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", content_types)
        archive.writestr("_rels/.rels", relationships)
        archive.writestr("word/document.xml", document)


def _write_xlsx(path: Path, marker: str) -> None:
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Synthetic"
    sheet["A1"] = marker
    workbook.save(path)


def _write_pptx(path: Path, marker: str) -> None:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    textbox = slide.shapes.add_textbox(914400, 914400, 7315200, 914400)
    textbox.text_frame.text = marker
    presentation.save(path)


def _write_png(path: Path, marker: str) -> None:
    from PIL import Image, ImageDraw, PngImagePlugin

    image = Image.new("RGB", (320, 96), "white")
    ImageDraw.Draw(image).text((12, 36), marker, fill="black")
    metadata = PngImagePlugin.PngInfo()
    metadata.add_text("SyntheticLabel", marker)
    image.save(path, pnginfo=metadata)


def _write_wav(path: Path) -> None:
    sample_rate = 8000
    frame_count = 1200
    frames = bytearray()
    for index in range(frame_count):
        sample = int(2000 * math.sin(2 * math.pi * 440 * index / sample_rate))
        frames.extend(struct.pack("<h", sample))
    with wave.open(str(path), "wb") as stream:
        stream.setnchannels(1)
        stream.setsampwidth(2)
        stream.setframerate(sample_rate)
        stream.writeframes(bytes(frames))


def _write_rar(path: Path) -> None:
    name = b"folder/metadata.txt"
    data = b"synthetic metadata payload"
    block_header = struct.Struct("<HBHH")
    file_header = struct.Struct("<LLBLLBBHL")
    signature = b"Rar!\x1a\x07\x00"
    main = block_header.pack(0x90CF, 0x73, 0, 13) + b"\0" * 6
    flags = 0x8000
    dos_time = ((2026 - 1980) << 25) + (7 << 21) + (26 << 16)
    fields = file_header.pack(
        len(data),
        len(data),
        2,
        zlib.crc32(data) & 0xFFFFFFFF,
        dos_time,
        20,
        0x30,
        len(name),
        0x20,
    ) + name
    header_size = block_header.size + len(fields)
    header_without_crc = block_header.pack(0, 0x74, flags, header_size) + fields
    header_crc = zlib.crc32(header_without_crc[2:]) & 0xFFFF
    header = block_header.pack(header_crc, 0x74, flags, header_size) + fields
    end_without_crc = block_header.pack(0, 0x7B, 0, block_header.size)
    end_crc = zlib.crc32(end_without_crc[2:]) & 0xFFFF
    end = block_header.pack(end_crc, 0x7B, 0, block_header.size)
    path.write_bytes(signature + main + header + data + end)


def _write_7z(path: Path) -> None:
    import py7zr

    with py7zr.SevenZipFile(path, "w") as archive:
        archive.writestr(b"synthetic metadata payload", "folder/metadata.txt")


def create_fixtures(root: str | Path) -> dict[str, Path]:
    target = Path(root)
    target.mkdir(parents=True, exist_ok=True)
    fixtures = {
        "txt": target / "synthetic.txt",
        "py": target / "synthetic.py",
        "pdf": target / "synthetic.pdf",
        "docx": target / "synthetic.docx",
        "xlsx": target / "synthetic.xlsx",
        "pptx": target / "synthetic.pptx",
        "png": target / "synthetic.png",
        "wav": target / "synthetic.wav",
        "mp4": target / "synthetic.mp4",
        "zip": target / "synthetic.zip",
        "rar": target / "synthetic.rar",
        "7z": target / "synthetic.7z",
        "exe": target / "synthetic.exe",
        "apk": target / "synthetic.apk",
    }
    fixtures["txt"].write_text(MARKERS["txt"], encoding="utf-8")
    fixtures["py"].write_text(
        f"VALUE = {MARKERS['py']!r}\n",
        encoding="utf-8",
    )
    _write_pdf(fixtures["pdf"], MARKERS["pdf"])
    _write_docx(fixtures["docx"], MARKERS["docx"])
    _write_xlsx(fixtures["xlsx"], MARKERS["xlsx"])
    _write_pptx(fixtures["pptx"], MARKERS["pptx"])
    _write_png(fixtures["png"], MARKERS["png"])
    _write_wav(fixtures["wav"])
    fixtures["mp4"].write_bytes(
        b"\x00\x00\x00\x18ftypisom\x00\x00\x02\x00isomiso2"
        + MARKERS["mp4"].encode("ascii")
    )
    with zipfile.ZipFile(fixtures["zip"], "w", zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("folder/readme.txt", MARKERS["zip"])
    _write_rar(fixtures["rar"])
    _write_7z(fixtures["7z"])
    fixtures["exe"].write_bytes(b"MZ" + b"synthetic-not-executable")
    with zipfile.ZipFile(fixtures["apk"], "w", zipfile.ZIP_STORED) as archive:
        archive.writestr(
            "AndroidManifest.xml",
            "<manifest package=\"invalid.synthetic.qa\" />",
        )
    return fixtures


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("output_dir")
    args = parser.parse_args()
    create_fixtures(args.output_dir)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
